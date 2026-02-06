from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Cria a apresentação
prs = Presentation()

# Função auxiliar para criar slides de texto
def criar_slide_conteudo(titulo, topicos):
    slide_layout = prs.slide_layouts[1] # Layout Título + Conteúdo
    slide = prs.slides.add_slide(slide_layout)

    # Define Título
    title = slide.shapes.title
    title.text = titulo

    # Adiciona os bullets
    tf = slide.placeholders[1].text_frame
    tf.text = topicos[0] # Primeiro tópico

    for topico in topicos[1:]:
        p = tf.add_paragraph()
        p.text = topico
        p.level = 0

# --- SLIDE 1: CAPA ---
slide_layout = prs.slide_layouts[0] # Layout de Capa
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "CITSM Analyzer"
subtitle.text = "Inteligência Artificial Aplicada à Gestão de Serviços de TI\n\nStack: Python + GPU Computing + Oracle"

# --- SLIDE 2: O PROBLEMA ---
criar_slide_conteudo(
    "1. O Desafio: Dados Não Estruturados",
    [
        "Volume alto de tickets mensais com texto livre.",
        "SQL tradicional falha em entender contextos ('Lento' vs 'Travando').",
        "Dificuldade em identificar duplicidades e re-trabalho.",
        "Necessidade de análise manual demorada para gerar indicadores."
    ]
)

# --- SLIDE 3: A SOLUÇÃO ---
criar_slide_conteudo(
    "2. A Solução: Ecossistema Inteligente",
    [
        "Dashboards Operacionais: Filtros em cascata (Data > Contrato > Serviço).",
        "Busca Semântica: Motor de busca que entende a 'intenção' do usuário.",
        "IA Local (On-Premise): Toda a inteligência roda na GPU local (RTX 3060).",
        "Privacidade Total: Nenhum dado sai da empresa para APIs externas."
    ]
)

# --- SLIDE 4: TECNOLOGIAS (DETALHADO) ---
# Vamos criar um slide com mais tópicos para caber tudo
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "3. Arquitetura Técnica & Stack"

# Vamos usar uma fonte menor para caber a lista completa
tf = slide.placeholders[1].text_frame
tf.clear() # Limpa o padrão

# Função para adicionar linhas
def add_line(texto, nivel=0):
    p = tf.add_paragraph()
    p.text = texto
    p.level = nivel
    p.font.size = Pt(20 if nivel == 0 else 16) # Tamanho da fonte ajustado

# Conteúdo Organizado por Camadas
add_line("Frontend & Visualização (Abas do Sistema):")
add_line("Streamlit: Framework de Web App Interativo.", 1)
add_line("Plotly: Renderização de gráficos dinâmicos e timelines.", 1)
add_line("Pandas: Manipulação de Dataframes em memória (ETL).", 1)

add_line("Núcleo de Inteligência Artificial (Backend):")
add_line("PyTorch + CUDA: Computação paralela na GPU (RTX 3060).", 1)
add_line("Sentence-Transformers: Gerenciamento de Embeddings (Modelo E5-Large).", 1)
add_line("BERTopic: Modelagem de tópicos e clusterização.", 1)

add_line("Processamento de Linguagem Natural (NLP):")
add_line("Scikit-Learn: Vetorização (CountVectorizer) e métricas.", 1)
add_line("NLTK: Tratamento de Stopwords e limpeza de texto em Português.", 1)

add_line("Infraestrutura de Dados:")
add_line("Oracle Database: Fonte da verdade (ODS_ITSM).", 1)

# --- SLIDE 5: RESULTADOS ---
criar_slide_conteudo(
    "4. Resultados Alcançados",
    [
        "Velocidade: Indexação de milhares de chamados em segundos.",
        "Assertividade: Busca encontra tickets mesmo sem palavras exatas.",
        "Auditoria: Detecção automática de chamados duplicados (>90% similaridade).",
        "Custo: Zero custo mensal de API (OpenAI/Azure)."
    ]
)

# --- SLIDE 6: DEMONSTRAÇÃO ---
criar_slide_conteudo(
    "5. Demonstração Prática",
    [
        "(Momento de alternar para o Sistema Real)",
        "",
        "1. Mostrar Filtros de Data e Contrato.",
        "2. Realizar uma Busca Semântica (Ex: 'Problema financeiro').",
        "3. Mostrar a detecção de Duplicados.",
        "4. Visualizar os Tópicos gerados pelo BERTopic."
    ]
)

# --- SLIDE 7: FUTURO ---
criar_slide_conteudo(
    "6. Próximos Passos (Roadmap)",
    [
        "QA de IA: Criação de 'Golden Set' para validar precisão.",
        "Feedback Loop: Botões de Like/Dislike para aprendizado contínuo.",
        "LLM Local: Implementação de Llama-3 para gerar resumos automáticos.",
        "Expansão: Aplicar o modelo para outros contratos/áreas."
    ]
)

# Salva o arquivo
nome_arquivo = 'Apresentacao_CITSM_Export.pptx'
prs.save(nome_arquivo)

print(f"✅ Arquivo '{nome_arquivo}' gerado com sucesso!")